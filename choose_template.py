# import os
# import openai
# import numpy as np
# from numpy.linalg import norm
# from pptx import Presentation

# openai.api_key = os.getenv("OPENAI_API_KEY")
# EMB_MODEL = "text-embedding-ada-002"
# THRESHOLD = 0.75

# def extract_text(template_path: str) -> str:
#     prs = Presentation(template_path)
#     texts = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text.strip():
#                 texts.append(shape.text.strip())
#     return "\n".join(texts)

# def get_embedding(text: str) -> np.ndarray:
#     resp = openai.Embedding.create(model=EMB_MODEL, input=text)
#     return np.array(resp["data"][0]["embedding"])

# def choose_template(prompt: str, templates_dir="templates/") -> str:
#     user_emb = get_embedding(prompt)
#     best_path = None
#     best_sim  = -1.0

#     for fname in os.listdir(templates_dir):
#         if not fname.lower().endswith(".pptx"):
#             continue
#         path = os.path.join(templates_dir, fname)
#         txt  = extract_text(path)
#         emb  = get_embedding(txt)
#         sim  = float(np.dot(user_emb, emb) / (norm(user_emb)*norm(emb)))
#         if sim > best_sim:
#             best_sim, best_path = sim, path

#     if best_sim < THRESHOLD:
#         return None
#     return best_path





# import os
# from dotenv import load_dotenv
# import numpy as np
# from numpy.linalg import norm
# from pptx import Presentation
# from openai import OpenAI

# # Load .env into environment
# load_dotenv()

# # Instantiate the OpenAI client
# client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# EMB_MODEL = "text-embedding-ada-002"
# THRESHOLD = 0.75

# def extract_text(template_path: str) -> str:
#     prs = Presentation(template_path)
#     texts = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text.strip():
#                 texts.append(shape.text.strip())
#     return "\n".join(texts)

# def get_embedding(text: str) -> np.ndarray:
#     resp = client.embeddings.create(model=EMB_MODEL, input=text)
#     return np.array(resp.data[0].embedding)

# def choose_template(prompt: str, templates_dir="templates/") -> str:
#     user_emb = get_embedding(prompt)
#     best_path = None
#     best_sim  = -1.0

#     for fname in os.listdir(templates_dir):
#         if not fname.lower().endswith(".pptx"):
#             continue
#         path = os.path.join(templates_dir, fname)
#         txt  = extract_text(path)
#         emb  = get_embedding(txt)
#         sim  = float(np.dot(user_emb, emb) / (norm(user_emb)*norm(emb)))
#         if sim > best_sim:
#             best_sim, best_path = sim, path

#     if best_sim < THRESHOLD:
#         return None
#     return best_path




import os
from dotenv import load_dotenv
import numpy as np
from numpy.linalg import norm
from pptx import Presentation
from openai import OpenAI
import random

# Load .env into environment
load_dotenv()

# Instantiate the OpenAI client
# client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

os.environ["OPENAI_API_KEY"] = st.secrets["OPENAI_API_KEY"]
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

EMB_MODEL = "text-embedding-ada-002"
THRESHOLD = 0.75

def extract_text(template_path: str) -> str:
    prs = Presentation(template_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)

def get_embedding(text: str) -> np.ndarray:
    resp = client.embeddings.create(model=EMB_MODEL, input=text)
    return np.array(resp.data[0].embedding)

def choose_template(prompt: str, templates_dir="templates/") -> str:
    user_emb = get_embedding(prompt)
    best_path = None
    best_sim  = -1.0
    all_templates = []

    # Scan through templates and collect them
    for fname in os.listdir(templates_dir):
        if not fname.lower().endswith(".pptx"):
            continue
        path = os.path.join(templates_dir, fname)
        all_templates.append(path)

        txt  = extract_text(path)
        emb  = get_embedding(txt)
        sim  = float(np.dot(user_emb, emb) / (norm(user_emb)*norm(emb)))
        
        if sim > best_sim:
            best_sim, best_path = sim, path

    # If no template matches the prompt well, pick a random one
    if best_sim < THRESHOLD:
        print(f"No matching template found above threshold {THRESHOLD}. Choosing a random template.")
        if all_templates:
            return random.choice(all_templates)
        else:
            raise FileNotFoundError(f"No PowerPoint templates found in directory: {templates_dir}")
    
    print(f"Found matching template with similarity {best_sim:.2f}")
    return best_path