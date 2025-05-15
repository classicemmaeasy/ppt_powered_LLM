# from pptx import Presentation

# def build_presentation(slide_data: dict, template_path: str, output_path="output.pptx") -> str:
#     prs = Presentation(template_path)
#     # pick a default layout: usually slide_layouts[1] is Title+Content
#     layouts = prs.slide_layouts
#     default_layout = layouts[1] if len(layouts) > 1 else layouts[0]

#     for s in slide_data["slides"]:
#         slide = prs.slides.add_slide(default_layout)
#         # Title
#         slide.shapes.title.text = s["title"]
#         # Body bullets
#         if len(slide.placeholders) > 1:
#             tf = slide.placeholders[1].text_frame
#             tf.clear()
#             for b in s["body"]:
#                 p = tf.add_paragraph()
#                 p.text = b

#     prs.save(output_path)
#     return output_path



# import os
# from dotenv import load_dotenv
# import numpy as np
# from numpy.linalg import norm
# from pptx import Presentation
# from pptx.enum.shapes import PP_PLACEHOLDER
# from openai import OpenAI
# from pptx.oxml.ns import qn

# load_dotenv()
# client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
# EMB_MODEL     = "text-embedding-ada-002"
# LAYOUT_THRESH = 0.60

# def get_embedding(text: str) -> np.ndarray:
#     resp = client.embeddings.create(model=EMB_MODEL, input=text)
#     return np.array(resp.data[0].embedding)

# def extract_layout_text(layout) -> str:
#     """Concatenate all placeholder default text or fallback to layout.name."""
#     texts = []
#     for ph in layout.placeholders:
#         if hasattr(ph, "text") and ph.text and ph.text.strip():
#             texts.append(ph.text.strip())
#     return "\n".join(texts) or layout.name

# def remove_all_slides(prs: Presentation):
#     """Remove every slide by clearing the root <p:sldIdLst>."""
#     root = prs._element  # p:presentation
#     sldIdLst = root.find(qn('p:sldIdLst'))
#     if sldIdLst is not None:
#         root.remove(sldIdLst)

# def build_presentation(slide_data: dict,
#                        template_path: str,
#                        output_path="output.pptx") -> str:
#     # 1. Load template and clear its slides
#     prs = Presentation(template_path)
#     remove_all_slides(prs)
#     assert len(prs.slides) == 0, "Template slides weren’t cleared!"

#     # 2. Gather “fillable” layouts: those with both TITLE & BODY
#     usable = []
#     for idx, layout in enumerate(prs.slide_layouts):
#         types = {ph.placeholder_format.type for ph in layout.placeholders}
#         if PP_PLACEHOLDER.TITLE in types and PP_PLACEHOLDER.BODY in types:
#             usable.append((idx, layout))
#     if not usable:
#         raise ValueError("No layouts with both TITLE+BODY placeholders found.")

#     layout_indices, layouts = zip(*usable)
#     layout_embs = [get_embedding(extract_layout_text(l)) for l in layouts]

#     # 3. For each generated slide, pick best layout & add it
#     for slide_info in slide_data["slides"]:
#         # embed combined title + bullets
#         content = slide_info["title"] + "\n" + "\n".join(slide_info.get("body", []))
#         emb = get_embedding(content)

#         sims = [float(np.dot(emb, le)/(norm(emb)*norm(le))) for le in layout_embs]
#         best_local = int(np.argmax(sims))
#         if sims[best_local] < LAYOUT_THRESH:
#             best_local = 0  # fallback
#         layout = prs.slide_layouts[layout_indices[best_local]]
#         slide = prs.slides.add_slide(layout)

#         # 4. Populate TITLE placeholder
#         for ph in slide.placeholders:
#             if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#                 ph.text = slide_info["title"]
#                 break

#         # 5. Populate BODY placeholder
#         populated = False
#         for ph in slide.placeholders:
#             if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
#                 tf = ph.text_frame
#                 tf.clear()
#                 for b in slide_info.get("body", []):
#                     p = tf.add_paragraph()
#                     p.text = b
#                 populated = True
#                 break

#         # 6. Final fallback: if somehow no BODY placeholder, append under TITLE
#         if not populated:
#             # find title placeholder again
#             for ph in slide.placeholders:
#                 if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#                     tf = ph.text_frame
#                     for b in slide_info.get("body", []):
#                         p = tf.add_paragraph()
#                         p.text = b
#                     break

#     # 7. Save
#     prs.save(output_path)
#     return output_path


# import os
# from dotenv import load_dotenv
# import numpy as np
# from numpy.linalg import norm
# from pptx import Presentation
# from openai import OpenAI
# from pptx.oxml.ns import qn

# # ─── Config & Helpers ───────────────────────────────────────────────────────────
# load_dotenv()
# client       = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
# EMB_MODEL    = "text-embedding-ada-002"
# LAYOUT_THRESH = 0.6  # tune to make layout choice stricter/looser

# def get_embedding(text: str) -> np.ndarray:
#     resp = client.embeddings.create(model=EMB_MODEL, input=text)
#     return np.array(resp.data[0].embedding)

# def extract_layout_text(layout) -> str:
#     """Concatenate all default placeholder texts or fallback to layout.name."""
#     texts = []
#     for ph in layout.placeholders:
#         if hasattr(ph, "text") and ph.text and ph.text.strip():
#             texts.append(ph.text.strip())
#     return "\n".join(texts) or layout.name

# def remove_all_slides(prs: Presentation):
#     """Remove all existing slides by clearing the <p:sldIdLst> element."""
#     root    = prs._element  # <p:presentation>
#     sldIdLst = root.find(qn('p:sldIdLst'))
#     if sldIdLst is not None:
#         root.remove(sldIdLst)

# # ─── Main Build Function ────────────────────────────────────────────────────────
# def build_presentation(slide_data: dict,
#                        template_path: str,
#                        output_path="output.pptx") -> str:
#     # 1) Load template & clear existing slides
#     prs = Presentation(template_path)
#     remove_all_slides(prs)
#     assert len(prs.slides) == 0, "Failed to clear template slides"

#     # 2) Find “fillable” layouts (at least one text placeholder) & embed them
#     usable = []
#     for idx, layout in enumerate(prs.slide_layouts):
#         # skip completely blank layouts (no placeholders)
#         if not layout.placeholders:
#             continue
#         usable.append((idx, layout))

#     if not usable:
#         raise ValueError("No layouts with placeholders found in template")

#     layout_indices, layouts = zip(*usable)
#     layout_embs = [get_embedding(extract_layout_text(l)) for l in layouts]

#     # 3) Iterate over each AI‑generated slide
#     for slide_info in slide_data["slides"]:
#         # embed title + bullets
#         content = slide_info["title"] + "\n" + "\n".join(slide_info.get("body", []))
#         emb     = get_embedding(content)

#         # pick best layout by cosine sim
#         sims      = [float(np.dot(emb, le)/(norm(emb)*norm(le))) for le in layout_embs]
#         best_loc  = int(np.argmax(sims))
#         if sims[best_loc] < LAYOUT_THRESH:
#             best_loc = 0  # fallback
#         layout = prs.slide_layouts[layout_indices[best_loc]]
#         slide  = prs.slides.add_slide(layout)

#         # 4) Gather all shapes that support text
#         text_shapes = [sh for sh in slide.shapes if getattr(sh, "has_text_frame", False)]

#         # 5) Title → first text shape
#         if text_shapes:
#             text_shapes[0].text = slide_info["title"]
#         else:
#             # as last resort: add a new textbox
#             left = top = width = height = None  # defaults: let pptx pick
#             txBox = slide.shapes.add_textbox(left, top, width, height)
#             text_shapes = [txBox]
#             text_shapes[0].text = slide_info["title"]

#         # 6) Body bullets → second text shape (if exists), else under title
#         bullets = slide_info.get("body", [])
#         if len(text_shapes) > 1:
#             tf = text_shapes[1].text_frame
#             tf.clear()
#             for b in bullets:
#                 p = tf.add_paragraph()
#                 p.text = b
#         else:
#             # append to the title shape’s frame
#             tf = text_shapes[0].text_frame
#             for b in bullets:
#                 p = tf.add_paragraph()
#                 p.text = b

#     # 7) Save and return
#     prs.save(output_path)
#     return output_path



# import os
# from dotenv import load_dotenv
# import numpy as np
# from numpy.linalg import norm
# from pptx import Presentation
# from pptx.enum.shapes import PP_PLACEHOLDER
# from pptx.util import Inches
# from openai import OpenAI
# from pptx.oxml.ns import qn

# # ─── Configuration ─────────────────────────────────────────────────────────────
# load_dotenv()
# client        = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
# EMB_MODEL     = "text-embedding-ada-002"
# LAYOUT_THRESH = 0.6  # tune as needed

# # ─── Helpers ───────────────────────────────────────────────────────────────────
# def get_embedding(text: str) -> np.ndarray:
#     resp = client.embeddings.create(model=EMB_MODEL, input=text)
#     return np.array(resp.data[0].embedding)

# def extract_layout_text(layout) -> str:
#     """Concatenate all placeholder texts or fallback to layout.name."""
#     texts = []
#     for ph in layout.placeholders:
#         if hasattr(ph, "text") and ph.text and ph.text.strip():
#             texts.append(ph.text.strip())
#     return "\n".join(texts) or layout.name

# def remove_all_slides(prs: Presentation):
#     """
#     Clears out all <p:sldId> children from the <p:sldIdLst> element,
#     but leaves the <p:sldIdLst> element itself intact.
#     """
#     root = prs.part._element  # this is the <p:presentation> XML element
#     sldIdLst = root.find(qn('p:sldIdLst'))
#     if sldIdLst is not None:
#         for child in list(sldIdLst):
#             sldIdLst.remove(child)

# # ─── Main Build Function ───────────────────────────────────────────────────────
# def build_presentation(slide_data: dict,
#                        template_path: str,
#                        output_path: str = "output.pptx") -> str:
#     # 1. Load the CBRE template and clear its existing slides
#     prs = Presentation(template_path)
#     remove_all_slides(prs)
#     assert len(prs.slides) == 0, "Failed to clear template slides!"

#     # 2. Identify fillable layouts (those with at least one placeholder)
#     usable = []
#     for idx, layout in enumerate(prs.slide_layouts):
#         if layout.placeholders:
#             usable.append((idx, layout))
#     if not usable:
#         raise ValueError("No layouts with placeholders found in the template.")

#     layout_indices, layouts = zip(*usable)
#     layout_embs = [get_embedding(extract_layout_text(l)) for l in layouts]

#     # 3. For each generated slide, pick the best layout & add it
#     for slide_info in slide_data["slides"]:
#         # embed combined title + bullets
#         content = slide_info["title"] + "\n" + "\n".join(slide_info.get("body", []))
#         emb     = get_embedding(content)

#         sims      = [float(np.dot(emb, le)/(norm(emb)*norm(le))) for le in layout_embs]
#         best_local = int(np.argmax(sims))
#         if sims[best_local] < LAYOUT_THRESH:
#             best_local = 0
#         layout = prs.slide_layouts[layout_indices[best_local]]
#         slide  = prs.slides.add_slide(layout)

#         # 4. Collect all text-capable shapes
#         text_shapes = [sh for sh in slide.shapes if getattr(sh, "has_text_frame", False)]

#         # 5. Put title into first text shape (or new textbox)
#         if text_shapes:
#             text_shapes[0].text = slide_info["title"]
#         else:
#             left   = Inches(1)
#             top    = Inches(1)
#             width  = prs.slide_width  - Inches(2)
#             height = prs.slide_height - Inches(2)
#             txBox  = slide.shapes.add_textbox(left, top, width, height)
#             text_shapes = [txBox]
#             text_shapes[0].text = slide_info["title"]

#         # 6. Put bullets into second text shape, or under title if only one exists
#         bullets = slide_info.get("body", [])
#         if len(text_shapes) > 1:
#             tf = text_shapes[1].text_frame
#             tf.clear()
#             for b in bullets:
#                 p = tf.add_paragraph()
#                 p.text = b
#         else:
#             tf = text_shapes[0].text_frame
#             for b in bullets:
#                 p = tf.add_paragraph()
#                 p.text = b

#     # 7. Save and return
#     prs.save(output_path)
#     return output_path


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def build_presentation(slide_data: dict, template_path: str, output_path="output.pptx") -> str:
    prs = Presentation(template_path)

    # Clear out all existing slides in the template
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    for slide_info in slide_data["slides"]:
        # Skip slide if image is specified in content
        if "image_path" in slide_info:
            continue

        # Try to pick a layout that has both title and at least one text frame
        selected_layout = None
        for layout in prs.slide_layouts:
            has_title = any(sh.is_placeholder and sh.placeholder_format.type == 1 for sh in layout.placeholders)
            has_body = any(sh.is_placeholder and sh.placeholder_format.type in (2, 3, 4) for sh in layout.placeholders)
            if has_title and has_body:
                selected_layout = layout
                break

        # Fallback if no good layout found
        if selected_layout is None:
            selected_layout = prs.slide_layouts[0]

        # Add slide with selected layout
        slide = prs.slides.add_slide(selected_layout)

        # Remove all image shapes
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide.shapes._spTree.remove(shape._element)
            elif shape.is_placeholder and shape.placeholder_format.type == 18:  # PICTURE placeholder
                slide.shapes._spTree.remove(shape._element)

        # Set title (if available)
        if slide.shapes.title:
            slide.shapes.title.text = slide_info.get("title", "")

        # Find the first usable body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx != 0 and shape.has_text_frame:
                tf = shape.text_frame
                tf.clear()
                for para in slide_info.get("body", []):
                    p = tf.add_paragraph()
                    p.text = para
                break  # only fill the first body-like placeholder

    prs.save(output_path)
    return output_path
